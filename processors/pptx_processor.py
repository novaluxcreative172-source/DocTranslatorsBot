"""
processors/pptx_processor.py
Translates a .pptx file: every shape's text frame (title, body, text
boxes), table cells, and speaker notes. Slide layout/design is untouched --
only the text content changes.
"""

from pptx import Presentation
from translator import translate_text


def _translate_text_frame(text_frame, target_lang):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text and run.text.strip():
                run.text = translate_text(run.text, target_lang)


def _translate_shape(shape, target_lang):
    if shape.has_text_frame:
        _translate_text_frame(shape.text_frame, target_lang)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _translate_text_frame(cell.text_frame, target_lang)

    # Group shapes contain nested shapes
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub_shape in shape.shapes:
            _translate_shape(sub_shape, target_lang)


def process(input_path: str, output_path: str, target_lang: str) -> None:
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            _translate_shape(shape, target_lang)

        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            _translate_text_frame(notes_tf, target_lang)

    prs.save(output_path)
